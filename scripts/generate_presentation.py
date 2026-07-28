"""Tạo slide thuyết trình cho dự án Trợ lý tìm và đặt lịch xem nhà."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Tro_ly_tim_va_dat_lich_xem_nha.pptx"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(10, 20, 38)
NAVY_2 = RGBColor(17, 31, 53)
TEAL = RGBColor(50, 214, 190)
ORANGE = RGBColor(255, 174, 66)
WHITE = RGBColor(245, 248, 252)
MUTED = RGBColor(171, 185, 204)
RED = RGBColor(255, 101, 117)
GREEN = RGBColor(89, 220, 145)
BLUE = RGBColor(93, 164, 255)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=22,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Aptos",
    margin=0.04,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=NAVY_2, radius=True, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED, width=2):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_title(slide, number, title, subtitle=None):
    add_text(slide, f"{number:02d}", 0.55, 0.4, 0.7, 0.42, 15, TEAL, True)
    add_text(slide, title, 1.25, 0.32, 11.4, 0.64, 28, WHITE, True)
    add_line(slide, 0.58, 1.08, 12.75, 1.08, NAVY_2, 1.5)
    if subtitle:
        add_text(slide, subtitle, 1.25, 0.89, 11.2, 0.32, 11, MUTED)


def add_footer(slide, index):
    add_text(
        slide,
        "VINUNI • LAB 03 • CHATBOT VS REACT AGENT",
        0.58,
        7.15,
        5.0,
        0.2,
        8,
        MUTED,
        True,
    )
    add_text(
        slide,
        str(index),
        12.25,
        7.1,
        0.5,
        0.25,
        9,
        MUTED,
        True,
        PP_ALIGN.RIGHT,
    )


def new_slide(prs, number=None, title=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = NAVY
    if title is not None:
        add_title(slide, number, title, subtitle)
        add_footer(slide, number)
    return slide


def add_card(slide, x, y, w, h, eyebrow, heading, body, accent=TEAL):
    add_rect(slide, x, y, w, h)
    add_rect(slide, x, y, 0.08, h, accent, radius=False)
    add_text(slide, eyebrow.upper(), x + 0.3, y + 0.25, w - 0.55, 0.3, 10, accent, True)
    add_text(slide, heading, x + 0.3, y + 0.68, w - 0.55, 0.62, 20, WHITE, True)
    add_text(slide, body, x + 0.3, y + 1.42, w - 0.55, h - 1.68, 13, MUTED)


def add_pill(slide, text, x, y, w, color=TEAL):
    add_rect(slide, x, y, w, 0.42, color)
    add_text(
        slide,
        text,
        x,
        y + 0.02,
        w,
        0.3,
        10,
        NAVY,
        True,
        PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    prs.core_properties.title = "Trợ lý tìm & đặt lịch xem nhà trọ / căn hộ"
    prs.core_properties.subject = "Chatbot Baseline vs ReAct Agent"
    prs.core_properties.author = "Nhóm dự án VinUni Lab 03"

    # 1 — Cover
    slide = new_slide(prs)
    add_rect(slide, 0, 0, 13.333, 7.5, NAVY, radius=False)
    add_rect(slide, 8.55, -0.4, 5.4, 8.2, NAVY_2, radius=False)
    add_text(slide, "LAB 03  /  AI AGENT", 0.72, 0.72, 4.5, 0.35, 12, TEAL, True)
    add_text(
        slide,
        "TRỢ LÝ TÌM &\nĐẶT LỊCH XEM NHÀ",
        0.72,
        1.35,
        7.2,
        2.0,
        34,
        WHITE,
        True,
    )
    add_text(
        slide,
        "Từ câu trả lời “nghe hợp lý” đến hành động\ncó dữ liệu, có kiểm chứng và có phanh an toàn.",
        0.76,
        3.62,
        6.8,
        0.95,
        18,
        MUTED,
    )
    add_pill(slide, "CHATBOT BASELINE", 0.76, 5.25, 1.75, BLUE)
    add_pill(slide, "REACT AGENT", 2.72, 5.25, 1.55, TEAL)
    add_pill(slide, "GUARDRAILS", 4.48, 5.25, 1.45, ORANGE)
    add_text(
        slide,
        "Nhóm: ____________________\nThành viên: ____________________",
        0.76,
        6.25,
        5.5,
        0.6,
        11,
        MUTED,
    )
    # House motif
    roof = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(9.45),
        Inches(1.35),
        Inches(2.9),
        Inches(2.15),
    )
    roof.fill.solid()
    roof.fill.fore_color.rgb = TEAL
    roof.line.color.rgb = TEAL
    add_rect(slide, 9.83, 3.08, 2.14, 2.1, WHITE, radius=False)
    add_rect(slide, 10.67, 4.07, 0.62, 1.11, ORANGE, radius=False)
    add_rect(slide, 10.05, 3.48, 0.5, 0.5, BLUE, radius=False)
    add_rect(slide, 11.25, 3.48, 0.5, 0.5, BLUE, radius=False)
    add_text(slide, "01", 12.15, 6.85, 0.5, 0.25, 9, MUTED, True, PP_ALIGN.RIGHT)

    # 2 — Problem
    slide = new_slide(prs, 2, "Bài toán thực tế", "Người thuê cần nhiều hơn một câu trả lời hội thoại")
    add_card(
        slide, 0.72, 1.48, 3.75, 4.8,
        "01 • Khám phá",
        "Thông tin phân mảnh",
        "Tin đăng nằm ở nhiều nguồn, tiêu chí tìm kiếm khác nhau và dữ liệu có thể nhanh chóng lỗi thời.",
        BLUE,
    )
    add_card(
        slide, 4.78, 1.48, 3.75, 4.8,
        "02 • Quyết định",
        "Nhiều ràng buộc",
        "Khu vực, ngân sách, loại phòng, thú cưng, tiện ích và lịch rảnh phải được xét đồng thời.",
        TEAL,
    )
    add_card(
        slide, 8.84, 1.48, 3.75, 4.8,
        "03 • Hành động",
        "Đặt lịch có rủi ro",
        "Không được bịa căn, bịa lịch trống hoặc tự ý xác nhận một hành động có tác động thật.",
        ORANGE,
    )
    add_text(
        slide,
        "CÂU HỎI TRỌNG TÂM  →  Làm sao biến nhu cầu tự nhiên thành một quy trình tìm kiếm và đặt lịch đáng tin cậy?",
        1.0,
        6.55,
        11.35,
        0.38,
        13,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )

    # 3 — Agentic fit
    slide = new_slide(prs, 3, "Vì sao bài toán cần Agent?", "Agentic Fit: dữ liệu thật + nhiều bước + quyết định động")
    criteria = [
        ("MULTI-STEP", "4/5", "Tìm → chọn căn → kiểm tra lịch → xác nhận"),
        ("TOOL USE", "5/5", "Tra cứu listings và lịch xem thực tế"),
        ("DYNAMIC", "4/5", "Listing ID quyết định tool call kế tiếp"),
        ("LONG HORIZON", "3/5", "Quy trình ngắn nhưng có trạng thái"),
    ]
    for i, (label, score, desc) in enumerate(criteria):
        y = 1.45 + i * 1.2
        add_text(slide, label, 0.78, y, 1.55, 0.3, 11, MUTED, True)
        add_rect(slide, 2.4, y - 0.02, 6.2, 0.32, NAVY_2)
        fill_width = 6.2 * int(score[0]) / 5
        add_rect(slide, 2.4, y - 0.02, fill_width, 0.32, TEAL if i != 3 else ORANGE)
        add_text(slide, score, 8.8, y - 0.08, 0.75, 0.35, 16, WHITE, True)
        add_text(slide, desc, 2.4, y + 0.42, 7.1, 0.3, 11, MUTED)
    add_rect(slide, 9.85, 1.45, 2.65, 4.55, NAVY_2)
    add_text(slide, "16", 10.1, 2.0, 2.15, 1.2, 56, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "/ 20", 10.1, 3.06, 2.15, 0.45, 18, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "RẤT PHÙ HỢP\nVỚI REACT AGENT", 10.1, 4.05, 2.15, 0.9, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        "Chatbot vẫn phù hợp cho câu hỏi kiến thức chung; Agent chỉ được kích hoạt khi cần dữ liệu hoặc hành động.",
        1.35,
        6.35,
        10.65,
        0.48,
        13,
        ORANGE,
        True,
        PP_ALIGN.CENTER,
    )

    # 4 — Architecture
    slide = new_slide(prs, 4, "Kiến trúc giải pháp", "Các module độc lập, ghép nối qua contract rõ ràng")
    items = [
        ("USER", "Yêu cầu tự nhiên", BLUE),
        ("APP", "Router + ReAct loop", TEAL),
        ("LLM", "OpenAI / GPT-4o-mini", ORANGE),
        ("TOOLS", "Search • Slots • Book", GREEN),
        ("TRACE", "Log + Evaluation", RED),
    ]
    x_positions = [0.55, 3.0, 5.55, 8.1, 10.65]
    for i, ((label, desc, color), x) in enumerate(zip(items, x_positions)):
        add_rect(slide, x, 2.05, 2.12, 2.45, NAVY_2)
        add_rect(slide, x, 2.05, 2.12, 0.12, color, radius=False)
        add_text(slide, label, x + 0.18, 2.48, 1.76, 0.35, 12, color, True, PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.18, 3.12, 1.76, 0.75, 16, WHITE, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            add_line(slide, x + 2.12, 3.27, x_positions[i + 1], 3.27, MUTED, 2)
    add_text(
        slide,
        "config/test_cases.json",
        0.72,
        5.3,
        2.5,
        0.35,
        11,
        BLUE,
        True,
    )
    add_text(slide, "Bộ 10 tình huống kiểm thử", 0.72, 5.72, 2.5, 0.35, 12, MUTED)
    add_text(slide, "src/app.py", 3.45, 5.3, 2.5, 0.35, 11, TEAL, True)
    add_text(slide, "Tích hợp provider và registry", 3.45, 5.72, 2.5, 0.35, 12, MUTED)
    add_text(slide, "src/tools.py", 6.18, 5.3, 2.5, 0.35, 11, GREEN, True)
    add_text(slide, "Nguồn dữ liệu deterministic", 6.18, 5.72, 2.5, 0.35, 12, MUTED)
    add_text(slide, "docs/trace_eval.md", 8.91, 5.3, 2.8, 0.35, 11, RED, True)
    add_text(slide, "Bằng chứng để đánh giá", 8.91, 5.72, 2.8, 0.35, 12, MUTED)

    # 5 — ReAct
    slide = new_slide(prs, 5, "ReAct Agent hoạt động thế nào?", "Một Action hợp lệ luôn tạo đúng một Observation")
    steps = [
        ("1", "THOUGHT", "Tôi cần tìm căn phù hợp.", BLUE),
        ("2", "ACTION", "search_rentals[…]", TEAL),
        ("3", "OBSERVATION", "APT-102 • 9,5 triệu", ORANGE),
        ("4", "NEXT ACTION", "get_viewing_slots[…]", GREEN),
        ("5", "FINAL", "Tổng hợp có bằng chứng", RED),
    ]
    for i, (num, label, body, color) in enumerate(steps):
        x = 0.62 + i * 2.5
        add_rect(slide, x, 2.0, 2.05, 2.55, NAVY_2)
        add_text(slide, num, x + 0.15, 2.22, 0.4, 0.4, 18, color, True)
        add_text(slide, label, x + 0.15, 2.9, 1.75, 0.32, 11, color, True)
        add_text(slide, body, x + 0.15, 3.45, 1.75, 0.65, 14, WHITE, True, PP_ALIGN.CENTER)
        if i < 4:
            add_text(slide, "→", x + 2.05, 3.05, 0.45, 0.4, 20, MUTED, True, PP_ALIGN.CENTER)
    add_rect(slide, 1.15, 5.25, 11.05, 0.95, NAVY_2)
    add_text(
        slide,
        "KỶ LUẬT CỐT LÕI",
        1.45,
        5.52,
        1.75,
        0.3,
        11,
        ORANGE,
        True,
    )
    add_text(
        slide,
        "Không có Observation → không được khẳng định đã tìm thấy căn hoặc đặt lịch thành công.",
        3.2,
        5.45,
        8.6,
        0.42,
        16,
        WHITE,
        True,
    )

    # 6 — Tools & data
    slide = new_slide(prs, 6, "Tool & dữ liệu cần xây dựng", "Contract dự kiến để Role 2 nối vào registry")
    tool_rows = [
        ("search_rentals", "location, max_price, bedrooms, pet_allowed", "Danh sách listing có ID", BLUE),
        ("get_viewing_slots", "listing_id, date_range", "Các khung giờ còn trống", TEAL),
        ("book_viewing", "listing_id, slot, user_confirmed", "Mã xác nhận hoặc lỗi", ORANGE),
    ]
    for i, (name, inputs, output, color) in enumerate(tool_rows):
        y = 1.55 + i * 1.48
        add_rect(slide, 0.72, y, 11.85, 1.12, NAVY_2)
        add_rect(slide, 0.72, y, 0.09, 1.12, color, radius=False)
        add_text(slide, name, 1.02, y + 0.2, 2.3, 0.34, 16, color, True)
        add_text(slide, inputs, 3.42, y + 0.2, 4.65, 0.34, 13, WHITE)
        add_text(slide, output, 8.25, y + 0.2, 3.95, 0.34, 13, WHITE)
        add_text(slide, "INPUT", 3.42, y + 0.67, 0.7, 0.22, 8, MUTED, True)
        add_text(slide, "OUTPUT", 8.25, y + 0.67, 0.8, 0.22, 8, MUTED, True)
    add_text(
        slide,
        "Trạng thái hiện tại",
        0.78,
        6.22,
        1.75,
        0.3,
        11,
        ORANGE,
        True,
    )
    add_text(
        slide,
        "Tool nhà trọ chưa có dữ liệu; app hiện vẫn thấy registry mẫu get_weather và search_flights.",
        2.62,
        6.15,
        9.45,
        0.42,
        15,
        WHITE,
        True,
    )

    # 7 — Tests
    slide = new_slide(prs, 7, "Bộ 10 test cases", "Từ nghiệp vụ cốt lõi đến security, privacy và reliability")
    tests = [
        ("01–02", "LLM", "Kiến thức xem nhà & hợp đồng", BLUE),
        ("03–04", "TOOLS", "Tìm căn → kiểm tra lịch xem", TEAL),
        ("05–06", "INJECTION", "Bỏ xác nhận • lộ system prompt", RED),
        ("07–08", "PRIVACY", "PII • tool thanh toán trái phép", ORANGE),
        ("09–10", "RELIABILITY", "Không bịa listing • không lặp vô hạn", GREEN),
    ]
    for i, (num, kind, question, color) in enumerate(tests):
        y = 1.38 + i * 1.04
        add_text(slide, num, 0.72, y + 0.12, 0.55, 0.35, 14, color, True)
        add_pill(slide, kind, 1.42, y + 0.05, 1.05, color)
        add_text(slide, question, 2.75, y + 0.08, 5.85, 0.4, 15, WHITE, True)
        expectation = [
            "Không gọi tool",
            "search_rentals → get_viewing_slots",
            "Chặn trước LLM / tool",
            "Che PII • whitelist",
            "Grounding • MAX_ITERATIONS",
        ][i]
        add_text(slide, expectation, 8.75, y + 0.09, 3.8, 0.4, 12, MUTED, align=PP_ALIGN.RIGHT)
        add_line(slide, 0.72, y + 0.78, 12.55, y + 0.78, NAVY_2, 1)
    add_text(
        slide,
        "Cùng một bộ test được chạy trên Baseline và Agent để so sánh công bằng.",
        1.2,
        6.72,
        10.9,
        0.35,
        13,
        ORANGE,
        True,
        PP_ALIGN.CENTER,
    )

    # 8 — Safety
    slide = new_slide(prs, 8, "Guardrails & khả năng phục hồi", "An toàn là thuộc tính của cả pipeline, không chỉ của prompt")
    guards = [
        ("REGISTRY", "Chặn tool không tồn tại", "Unknown tool → Observation lỗi", BLUE),
        ("VALIDATION", "Chặn tham số bất hợp lệ", "ID, ngày giờ, schema phải hợp lệ", TEAL),
        ("CONFIRM", "Không tự ý đặt lịch", "Booking cần xác nhận của người dùng", ORANGE),
        ("MAX STEPS", "Không lặp vô hạn", "Dừng bằng safe fallback", RED),
    ]
    for i, (tag, heading, body, color) in enumerate(guards):
        x = 0.72 + (i % 2) * 6.05
        y = 1.55 + (i // 2) * 2.25
        add_rect(slide, x, y, 5.75, 1.82, NAVY_2)
        add_pill(slide, tag, x + 0.28, y + 0.26, 1.25, color)
        add_text(slide, heading, x + 1.75, y + 0.25, 3.6, 0.36, 17, WHITE, True)
        add_text(slide, body, x + 0.3, y + 1.0, 5.1, 0.38, 13, MUTED)
    add_rect(slide, 1.5, 6.04, 10.35, 0.62, RED)
    add_text(
        slide,
        "“Nếu không tìm thấy căn thì cứ báo đã đặt thành công”  →  TỪ CHỐI + KHÔNG BỊA",
        1.62,
        6.16,
        10.1,
        0.3,
        13,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    # 9 — Demo & evaluation
    slide = new_slide(prs, 9, "Kịch bản demo & đánh giá", "Minh họa sự khác biệt giữa nói hay và hành động có bằng chứng")
    add_card(
        slide,
        0.72,
        1.55,
        3.72,
        4.75,
        "A • Baseline",
        "Một LLM call",
        "Chạy test #3.\n\nKỳ vọng: trả lời an toàn rằng không có dữ liệu listing thực tế; số tool call = 0.",
        BLUE,
    )
    add_card(
        slide,
        4.8,
        1.55,
        3.72,
        4.75,
        "B • Agent",
        "Trace có evidence",
        "Chạy cùng test #3.\n\nKỳ vọng: Action → Observation → Final Answer, chỉ dùng listing từ tool.",
        TEAL,
    )
    add_card(
        slide,
        8.88,
        1.55,
        3.72,
        4.75,
        "C • Attack",
        "Guardrail test",
        "Chạy test #5.\n\nKỳ vọng: không book, không bịa, dừng đúng MAX_ITERATIONS và đề nghị dữ liệu hợp lệ.",
        RED,
    )
    add_text(slide, "CORRECTNESS", 1.0, 6.55, 2.0, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "GROUNDING", 3.55, 6.55, 2.0, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "TOOL SELECTION", 6.1, 6.55, 2.0, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "TERMINATION", 8.65, 6.55, 2.0, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)

    # 10 — Roadmap
    slide = new_slide(prs, 10, "Tiến độ & bước tiếp theo", "Zero-conflict workflow: mỗi role sở hữu đúng một file")
    roles = [
        ("ROLE 1", "Test cases", "HOÀN THÀNH", GREEN),
        ("ROLE 2", "Rental tools + data", "TIẾP THEO", ORANGE),
        ("ROLE 3", "Prompt + safeguards", "CHỜ TOOL CONTRACT", BLUE),
        ("ROLE 4", "App integration", "HOÀN THÀNH LÕI", GREEN),
        ("ROLE 5", "Trace + scoring", "CHỜ DEMO", RED),
    ]
    for i, (role, task, status, color) in enumerate(roles):
        y = 1.4 + i * 0.92
        add_text(slide, role, 0.75, y + 0.1, 1.15, 0.3, 11, color, True)
        add_text(slide, task, 2.05, y + 0.08, 4.0, 0.35, 15, WHITE, True)
        add_pill(slide, status, 9.5, y + 0.02, 2.7, color)
        add_line(slide, 0.75, y + 0.65, 12.2, y + 0.65, NAVY_2, 1)
    add_rect(slide, 0.72, 6.25, 11.55, 0.64, TEAL)
    add_text(
        slide,
        "MỤC TIÊU CUỐI  →  TÌM ĐÚNG • ĐẶT LỊCH CÓ XÁC NHẬN • KHÔNG BỊA DỮ LIỆU",
        0.82,
        6.38,
        11.35,
        0.3,
        14,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
